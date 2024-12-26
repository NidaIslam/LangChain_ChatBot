
import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.llms import HuggingFaceHub
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

CSS = '''
<style>
.chat-message {
padding: 1.5rem; border-radius: 0.5rem; margin-bottom:1rem; display: flex
}
.chat-message.user{
    background-color: #2b313e    
}
.chat-message.bot{
    background-color: #233861
}
.chat-message .avatar {
width: 15%;
}
.chat-message .avatar img {
max-width: 78px;
max-height: 78px;
border-radius: 50%;
object-fit: cover;
}
.chat-message .message {
width: 85%;
padding: 0 1.5rem;
color: #fff;
}
'''
bot_template = '''
<div class="chat-message bot">
    <div class="avatar">
        <img src="https://cbx-prod.b-cdn.net/COLOURBOX30822090.jpg?width=800&height=800&quality=70" style="max-height: 50px; max-width: 50px; border-radius: 10%; object-fit: cover;">
    </div>
    <div class="message">{{MSG}}</div>
</div> 
'''
user_template = '''
<div class="chat-message user">
    <div class="avatar">
        <img src="https://cbx-prod.b-cdn.net/COLOURBOX3058324.jpg?width=800&height=800&quality=70" style="max-height: 50px; max-width: 50px; border-radius: 10%; object-fit: cover;">
    </div>
    <div class="message">{{MSG}}</div>
</div>
'''


def get_file_text(uploaded_files):
    """
    Extracts text from PDF, TXT, Word (DOCX), and PowerPoint (PPTX) files.
    """
    text = ""
    try:
        for uploaded_file in uploaded_files:
            file_name = uploaded_file.name.lower()
            
            # Check if file is a PDF
            if file_name.endswith(".pdf"):
                pdf_reader = PdfReader(uploaded_file)
                for page in pdf_reader.pages:
                    text += page.extract_text() or ""
            
            # Check if file is a TXT file
            elif file_name.endswith(".txt"):
                text += uploaded_file.read().decode("utf-8")
            
            # Check if file is a Word document
            elif file_name.endswith(".docx"):
                doc = Document(uploaded_file)
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
            
            # Check if file is a PowerPoint presentation
            elif file_name.endswith(".pptx"):
                ppt = Presentation(uploaded_file)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            
            # Unsupported file type
            else:
                raise ValueError(f"Unsupported file type: {file_name}")
    
    except Exception as e:
        st.error(f"Error reading files: {e}")
    
    return text


def get_text_chunks(text):
    """
    Splits text into smaller chunks.
    """
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    """
    Creates a FAISS vector store using HuggingFace embeddings.
    """

    # Initialize the HuggingFace embedding model
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    
    # Create the FAISS vector store
    vectorstore = FAISS.from_texts(text_chunks, embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    """
    Initializes a conversational chain with a memory buffer.
    """
   
    llm = HuggingFaceHub(repo_id="google/flan-t5-large", 
        model_kwargs={
            "temperature": 0.4,
            "max_length": 500, 
            "top_p": 0.8,
            "repetition_penalty": 1.2
        }
    )
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    
    conversational_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversational_chain


def handle_userinput(user_question):
    """
    Handles user input and displays conversation.
    """
    if not st.session_state.conversation:
        st.error("Please upload and process your documents first.")
        return
    
    try:
        response = st.session_state.conversation({'question': user_question})
        st.session_state.chat_history = response['chat_history']
        for i, message in enumerate(st.session_state.chat_history):
            if i % 2 == 0:
                st.write(user_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
            else:
                st.write(bot_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
    except Exception as e:
        st.error(f"Error processing your question: {e}")


def main():
    """
    Main Streamlit app function.
    """

    load_dotenv()
    st.set_page_config(page_title="Chat with multiple Docs", page_icon=":books:")
    st.write(CSS, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    st.header("Chat with multiple Documents :robot_face::books:")

    with st.sidebar:
        st.subheader(":open_file_folder: Your documents")
        pdf_docs = st.file_uploader(
            "Upload your Documents here and click 'Process'", accept_multiple_files=True)
        if st.button("Process"):
            if pdf_docs:
                with st.spinner("Processing..."):
                    raw_text = get_file_text(pdf_docs)
                    text_chunks = get_text_chunks(raw_text)
                    vectorstore = get_vectorstore(text_chunks)
                    st.session_state.conversation = get_conversation_chain(vectorstore)
                st.success("Documents processed successfully!")
            else:
                st.error("Please upload at least one File.")

    user_question = st.text_input("Ask a question about your document: :mag:")
    if user_question:
        handle_userinput(user_question)

if __name__ == '__main__':
    main()
